import logging
from pdfplumber import PDF
from pptx import Presentation

from data_handler.data_helper import DataHelper

logger = logging.getLogger(__name__)
logging.basicConfig(level=logging.INFO)


class DataProcessor:
    @staticmethod
    def clean_json(data: dict):
        companies = []
        employees = []

        for company in data["companies"]:
            company_info = {
                "id": DataHelper.format_data(company.get("id")),
                "name": DataHelper.format_data(company.get("name")),
                "industry": DataHelper.format_data(company.get("industry")),
                "revenue": DataHelper.format_data(company.get("revenue")),
                "location": DataHelper.format_data(company.get("location")),
                "performance_2023_Q1_revenue": DataHelper.format_data(
                    company.get("performance", {}).get("2023_Q1", {}).get("revenue")
                ),
                "performance_2023_Q1_profit_margin": DataHelper.format_data(
                    company.get("performance", {})
                    .get("2023_Q1", {})
                    .get("profit_margin")
                ),
                "performance_2023_Q2_revenue": DataHelper.format_data(
                    company.get("performance", {}).get("2023_Q2", {}).get("revenue")
                ),
                "performance_2023_Q2_profit_margin": DataHelper.format_data(
                    company.get("performance", {})
                    .get("2023_Q2", {})
                    .get("profit_margin")
                ),
            }
            print(
                DataHelper.format_data(
                    company.get("performance", {}).get("2023_Q1", {}).get("revenue")
                )
            )
            companies.append(company_info)

            for emp in company["employees"]:
                emp_info = {
                    "company_name": DataHelper.format_data(company.get("name")),
                    "id": DataHelper.format_data(emp.get("id")),
                    "name": DataHelper.format_data(emp.get("name")),
                    "role": DataHelper.format_data(emp.get("role")),
                    "salary": DataHelper.format_data(emp.get("salary")),
                    "hired_date": DataHelper.format_date(emp.get("hired_date")),
                }
                employees.append(emp_info)
        cleaned_data = {"companies": companies, "employees": employees}
        return cleaned_data

    @staticmethod
    def clean_pptx(presentation: Presentation):
        data = {}
        for slide_num, slide in enumerate(presentation.slides):
            slide_data = {}

            text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text.strip())

            slide_data["text"] = "\n".join(text)

            table_data = []
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    headers = [cell.text for cell in table.rows[0].cells]

                    for row_index, row in enumerate(table.rows):
                        if row_index == 0:
                            continue
                        row_data = {
                            headers[i]: DataHelper.format_data(
                                row.cells[i].text.strip()
                            )
                            for i in range(len(headers))
                        }
                        table_data.append(row_data)

            if table_data:
                slide_data["table"] = table_data

            data[f"page_{slide_num + 1}"] = slide_data
        return data

    @staticmethod
    def clean_pdf(pdf: PDF):
        data = {}
        for page_num, page in enumerate(pdf.pages):
            page_data = {}

            text = page.extract_text()
            cleaned_text = text.split("\n")[0] if text else ""
            page_data["text"] = cleaned_text

            table = page.extract_table()
            if table:
                headers = table[0]
                table_data = []
                for row in table[1:]:
                    row_data = {
                        headers[i]: DataHelper.format_data(row[i])
                        for i in range(len(headers))
                    }
                    table_data.append(row_data)
                page_data["table"] = table_data
            else:
                page_data["table"] = None

            data[f"page_{page_num + 1}"] = page_data
        return data

    @staticmethod
    def clean_csv(datas: dict):
        clean_data = []
        for data in datas:
            info = {
                "date": DataHelper.format_date(data["Date"]),
                "membership_ID": DataHelper.format_data(data["Membership_ID"]),
                "membership_Type": DataHelper.format_data(data["Membership_Type"]),
                "activity": DataHelper.format_data(data["Activity"]),
                "revenue": DataHelper.format_data(data["Revenue"]),
                "duration(minutes)": DataHelper.format_data(data["Duration (Minutes)"]),
                "location": DataHelper.format_data(data["Location"]),
            }
            clean_data.append(info)
        return clean_data
